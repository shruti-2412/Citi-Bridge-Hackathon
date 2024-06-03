import streamlit as st
import PyPDF2
import json
import app
import docx
import pptx
import textract

def extract_text_from_file(uploaded_file):
    text = ''
    try:
        file_extension = uploaded_file.name.split('.')[-1]
        
        if file_extension == "pdf":
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page_num in range(len(pdf_reader.pages)):
                text += pdf_reader.pages[page_num].extract_text()
        elif file_extension == "docx":
            doc = docx.Document(uploaded_file)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        elif file_extension == "txt":
            text = uploaded_file.read().decode("utf-8")
        elif file_extension == "pptx":
            ppt = pptx.Presentation(uploaded_file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            # Using textract for other file types like doc, xls, etc.
            text = textract.process(uploaded_file).decode("utf-8")
            
    except Exception as e:
        st.error(f"Error extracting text from {file_extension.upper()} file: {e}")
    return text

def extract_text_from_json(uploaded_json):
    try:
        print(uploaded_json.name)
        # Read the input JSON file
        with open("text-summarizer/abc.json") as f:
            data = json.load(f)

        # Create a dictionary to store body texts for each thread_id
        thread_data = {}

        # Iterate over each message and organize by thread_id
        for message in data:
            thread_id = message['thread_id']
            body_text = message['body']

            if thread_id not in thread_data:
                thread_data[thread_id] = []

            thread_data[thread_id].append(body_text)

        # Create a list of dictionaries with thread_id and body texts
        output_data = []
        for thread_id, body_texts in thread_data.items():
            output_data.append({
                "thread_id": thread_id,
                "body": body_texts
                
            })

        return output_data

    except Exception as e:
        st.error(f"Error extracting JSON: {e}")
        return

def write_text_to_file(text, file_format):
    try:
        if file_format == "pdf":
            with open("summary.pdf", 'wb') as file:
                file.write(text.encode('utf-8'))
            st.write("PDF File written successfully")
        elif file_format == "docx":
            doc = docx.Document()
            doc.add_paragraph(text)
            doc.save("summary.docx")
            st.write("DOCX File written successfully")
        else:
            with open("summary.txt", 'w', encoding='utf-8') as file:
                file.write(text)
            st.write("TXT File written successfully")
    except Exception as e:
        st.error(f"Error writing to file: {e}")

def main():
    st.title("Document Text Summarizer")

    # Radio button to select input method
    button_choice = st.radio("Select Input Method", ("Upload File", "Use Text Boxes","Upload JSON"))

    if button_choice == "Upload File":
        # File upload button
        uploaded_file = st.file_uploader("Choose a file", type=["pdf", "docx", "txt", "pptx"])

        if uploaded_file is not None:
            # Display file upload confirmation
            st.write("File Uploaded Successfully!")

            # Option for summarization
            option = st.radio("Select the words limit of the summary", ("25%", "50%", "75%", "egular"))

            # Select download format
            download_format = st.selectbox("Select download format", ("DOCX", "TXT"))

            # Submit button
            if st.button("Summarize Text"):
                # Perform text extraction on file submission
                extracted_text = extract_text_from_file(uploaded_file)
                if extracted_text:
                    st.success("Wait for a while we are processing the your file......")

                    # Select summarization based on chosen option
                    if option == "25%":
                        summarytext = app.twentyfive(extracted_text)
                    elif option == "50%":
                        summarytext = app.fifty(extracted_text)
                    elif option == "75%":
                        summarytext = app.seventyfive(extracted_text)
                    else:
                        summarytext = app.generateSummary(extracted_text)
                    
                    st.write(summarytext)  # Display the summary text
                    if summarytext is not None:
                        write_text_to_file(summarytext, download_format.lower())
                        with open(f"summary.{download_format.lower()}", "rb") as file:
                            file_contents = file.read()
                        st.download_button(label=f"Download {download_format.upper()}", data=file_contents, file_name=f"summary.{download_format.lower()}")
    # for text to text
    elif button_choice == "Use Text Boxes":
          #st.title("Two Text Boxes Web Page")

        # Create two columns
        col1, col2 = st.columns(2)

        # Add text boxes to the columns
        with col1:
            left_text = st.text_area("Input Text", "Type something here", height=400)

        with col2:
        # Display the initial summary of text from the left text box
            right_text = st.empty()  # Placeholder for the right text box
            initial_summary = app.generateSummary(left_text)
            right_text.text_area("Summarized Text", initial_summary, disabled=True, height=400)

    # Add a submit button
        if st.button("Submit"):
        # Update the content of the right text box with the summary
            updated_summary = app.generateSummary(left_text)
            right_text.empty()  # Clear the placeholder
            right_text.text_area("Summarized Text", updated_summary, height=400, key="summarized_text")

    elif button_choice == "Upload JSON":
        # File upload button for JSON
        uploaded_json = st.file_uploader("Choose a JSON file", type=["json"])
       
        data = []  # This list will store dictionaries with thread ID and body
        if uploaded_json is not None:
            # Display file upload confirmation
            st.write("JSON File Uploaded Successfully!")

            if st.button("Summarize Text"):
                # Perform text extraction on PDF file submission
                extracted_json_data = extract_text_from_json(uploaded_json)
                if extracted_json_data:
                    st.success("Wait for a while we are processing your json file.....")

                summarytext = app.generateSummaryJSON(extracted_json_data)
                st.write(summarytext)

                with open("threads.json", "w") as outfile:
                    json.dump(summarytext, outfile, indent=4)  # Add indentation for readability (optional)

                print("JSON file created successfully!")
                if summarytext is not None:
                    
                    download_button = st.download_button(label="Download JSON file ", data=json.dumps(summarytext).encode("utf-8"), file_name="jsonSummary.json")
                    if download_button:
                        st.info("Text File Downloaded Successfully!")
            

if __name__ == "__main__":
    main()
